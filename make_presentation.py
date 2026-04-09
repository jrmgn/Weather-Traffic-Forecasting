# make_presentation.py
from pptx import Presentation
from pptx.util import Inches
import os

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Delivery Delay Forecasts vs Weather Trends"
slide.placeholders[1].text = "Smart Logistics Summary - Jermagne"

slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Delays vs Precipitation"
if os.path.exists("delays_vs_weather.png"):
    prs.slides[1].shapes.add_picture("delays_vs_weather.png", Inches(1), Inches(1.5), width=Inches(8))

prs.save("smart_logistics_summary.pptx")
print("Saved smart_logistics_summary.pptx")